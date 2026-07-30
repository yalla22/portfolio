# -*- coding: utf-8 -*-
"""Marut Drones — all 10 projects deck for the Yali Aerospace technical interview.
Each project slide has a screenshot placeholder on the right; drop images in later."""
import os
from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

NAVY = RGBColor(0x14, 0x2A, 0x47)
NAVY2 = RGBColor(0x1F, 0x3A, 0x5F)
ACCENT = RGBColor(0x2E, 0x9B, 0xDF)
TEXT = RGBColor(0x1E, 0x24, 0x2E)
MUTED = RGBColor(0x55, 0x5F, 0x6E)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT = RGBColor(0xC9, 0xD6, 0xE5)
PH_BG = RGBColor(0xEC, 0xEF, 0xF3)
PH_LINE = RGBColor(0xB6, 0xC2, 0xD2)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
SW, SH = prs.slide_width, prs.slide_height
IMG_DIR = r"c:\Users\yalla\Videos\DISC_14\ppt_images"  # drop screenshots here as 1.png, 2.png, ...


def _rect(slide, x, y, w, h, color, line_color=None, rounded=False):
    shape = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
    shp = slide.shapes.add_shape(shape, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    if line_color is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line_color
        shp.line.width = Pt(1)
    shp.shadow.inherit = False
    return shp


def _text(slide, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    for i, (txt, size, color, bold, sb, level) in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.level = level
        if sb:
            p.space_before = Pt(sb)
        p.space_after = Pt(3)
        r = p.add_run()
        r.text = txt
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color
        r.font.name = "Calibri"
    return tb


def _footer(s, n):
    _text(s, Inches(0.7), Inches(7.05), Inches(9), Inches(0.35),
          [("Janadri Yalla Yashwanth  ·  Marut Drones", 10, MUTED, False, 0, 0)])
    _text(s, Inches(12.3), Inches(7.05), Inches(0.7), Inches(0.35),
          [(str(n), 10, MUTED, False, 0, 0)], align=PP_ALIGN.RIGHT)


def title_slide(title, subtitle, presenter, meta):
    s = prs.slides.add_slide(BLANK)                    # white background
    _rect(s, 0, 0, Inches(0.26), SH, ACCENT)           # slim left accent strip
    _rect(s, Inches(0.95), Inches(2.5), Inches(1.3), Inches(0.08), ACCENT)
    _text(s, Inches(0.95), Inches(2.65), Inches(11.4), Inches(2.1),
          [(title, 44, NAVY2, True, 0, 0), (subtitle, 20, MUTED, False, 12, 0)])
    _text(s, Inches(0.95), Inches(5.9), Inches(11.4), Inches(1.2),
          [(presenter, 20, NAVY2, True, 0, 0), (meta, 14, MUTED, False, 4, 0)])


def section_slide(text, n):
    s = prs.slides.add_slide(BLANK)                    # white background
    _rect(s, Inches(0.95), Inches(3.25), Inches(1.2), Inches(0.08), ACCENT)
    _text(s, Inches(0.95), Inches(3.45), Inches(11.4), Inches(1.2), [(text, 32, NAVY2, True, 0, 0)])


def _header(s, kicker, title):
    # light, clean header: navy title on white with an accent underline
    _text(s, Inches(0.7), Inches(0.42), Inches(12), Inches(0.32), [(kicker, 12.5, ACCENT, True, 0, 0)])
    _text(s, Inches(0.7), Inches(0.7), Inches(12), Inches(0.65), [(title, 26, NAVY2, True, 0, 0)])
    _rect(s, Inches(0.7), Inches(1.36), Inches(11.93), Inches(0.014), RGBColor(0xDC, 0xE1, 0xE7))
    _rect(s, Inches(0.7), Inches(1.33), Inches(1.5), Inches(0.05), ACCENT)


def plain_slide(kicker, title, body, n):
    s = prs.slides.add_slide(BLANK)
    _header(s, kicker, title)
    runs = []
    for j, (txt, level) in enumerate(body):
        if level == 0:
            runs.append(("▸  " + txt, 19, TEXT, True, 10 if j else 6, 0))
        elif level == 1:
            runs.append(("•  " + txt, 16, MUTED, False, 3, 1))
        else:
            runs.append((txt, 14, ACCENT, True, 8, 0))
    _text(s, Inches(0.75), Inches(1.7), Inches(11.9), Inches(5.2), runs)
    _footer(s, n)


def project_slide(kicker, title, bullets, tech, n, img):
    s = prs.slides.add_slide(BLANK)
    _header(s, kicker, title)
    # left: bullets
    runs = []
    for j, (txt, level) in enumerate(bullets):
        if level == 0:
            runs.append(("▸  " + txt, 17, TEXT, True, 9 if j else 4, 0))
        else:
            runs.append(("•  " + txt, 14.5, MUTED, False, 3, 1))
    runs.append((tech, 13, ACCENT, True, 12, 0))
    has_img = bool(img and os.path.exists(img))
    _text(s, Inches(0.7), Inches(1.65), Inches(7.2) if has_img else Inches(11.9), Inches(5.2), runs)
    # right: image only if one exists (no placeholder box when absent)
    if has_img:
        fit_image(s, img, Inches(8.15), Inches(1.75), Inches(4.55), Inches(4.7))
    _footer(s, n)


REPORT_PIX = r"c:\Users\yalla\Videos\DISC_14\reports\Pixhawk_ArduPilot_report.pdf"
REPORT_JIYI = r"c:\Users\yalla\Videos\DISC_14\reports\JIYI_report.pdf"


def link_button(slide, x, y, w, h, label, target):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = ACCENT
    shp.line.fill.background()
    shp.shadow.inherit = False
    tf = shp.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = label
    r.font.size = Pt(13)
    r.font.bold = True
    r.font.color.rgb = WHITE
    r.font.name = "Calibri"
    shp.click_action.hyperlink.address = target
    return shp


def log_analyzer_slide(kicker, title, bullets, tech, n):
    s = prs.slides.add_slide(BLANK)
    _header(s, kicker, title)
    runs = []
    for j, (txt, level) in enumerate(bullets):
        if level == 0:
            runs.append(("▸  " + txt, 16, TEXT, True, 8 if j else 4, 0))
        else:
            runs.append(("•  " + txt, 13.5, MUTED, False, 3, 1))
    runs.append((tech, 12.5, ACCENT, True, 10, 0))
    _text(s, Inches(0.7), Inches(1.6), Inches(6.7), Inches(5.3), runs)
    # right panel — format detection by header magic bytes
    px, pw = Inches(7.75), Inches(4.9)
    _rect(s, px, Inches(1.65), pw, Inches(2.7), PH_BG, line_color=PH_LINE, rounded=True)
    _text(s, px + Inches(0.22), Inches(1.82), pw - Inches(0.44), Inches(2.4), [
        ("Format detection — header magic bytes", 14, NAVY2, True, 0, 0),
        ("Pixhawk / ArduPilot   →   A3 95", 13, TEXT, False, 9, 0),
        ("JIYI (plain)   →   00 FF", 13, TEXT, False, 4, 0),
        ("JIYI (encrypted)   →   37 C8", 13, TEXT, False, 4, 0),
        ("XOR every byte with 0x37 to decrypt", 12.5, ACCENT, True, 9, 0),
        ("37^37 = 00 · C8^37 = FF → the known JIYI marker", 11, MUTED, False, 3, 0),
    ])
    link_button(s, px, Inches(4.7), pw, Inches(0.72), "Open sample report — Pixhawk (ArduPilot)", REPORT_PIX)
    link_button(s, px, Inches(5.6), pw, Inches(0.72), "Open sample report — JIYI", REPORT_JIYI)
    _footer(s, n)


def fit_image(slide, path, bx, by, bw, bh):
    """Place an image scaled to fit inside the box (bx,by,bw,bh) EMU, centered, aspect preserved."""
    with Image.open(path) as im:
        iw, ih = im.size
    scale = min(bw / iw, bh / ih)
    w, h = int(iw * scale), int(ih * scale)
    slide.shapes.add_picture(path, bx + (bw - w) // 2, by + (bh - h) // 2, width=w, height=h)


def multi_img_slide(kicker, title, bullets, tech, n, imgs):
    s = prs.slides.add_slide(BLANK)
    _header(s, kicker, title)
    runs = []
    for j, (txt, level) in enumerate(bullets):
        if level == 0:
            runs.append(("▸  " + txt, 15, TEXT, True, 8 if j else 4, 0))
        else:
            runs.append(("•  " + txt, 12.5, MUTED, False, 3, 1))
    runs.append((tech, 12, ACCENT, True, 10, 0))
    _text(s, Inches(0.55), Inches(1.6), Inches(4.5), Inches(5.4), runs)
    have = [p for p in imgs if os.path.exists(p)]
    RX = Inches(5.25)
    if len(have) == 1:
        fit_image(s, have[0], RX, Inches(1.7), Inches(7.5), Inches(4.7))
    elif len(have) == 2:
        fit_image(s, have[0], RX, Inches(1.7), Inches(3.65), Inches(4.6))
        fit_image(s, have[1], RX + Inches(3.85), Inches(1.7), Inches(3.65), Inches(4.6))
    elif len(have) >= 3:
        fit_image(s, have[0], RX, Inches(1.6), Inches(3.75), Inches(2.5))
        fit_image(s, have[1], RX + Inches(3.9), Inches(1.6), Inches(3.75), Inches(2.5))
        fit_image(s, have[2], RX, Inches(4.25), Inches(7.65), Inches(2.45))
    _footer(s, n)


# map project title -> project number for placeholder label
_PROJ_NUM = {}


def n_of(title):
    return _PROJ_NUM.get(title, 0)


def img_for(i):
    for ext in (".png", ".jpg", ".jpeg"):
        p = os.path.join(IMG_DIR, "%d%s" % (i, ext))
        if os.path.exists(p):
            return p
    return None


# ---------------- BUILD ----------------
title_slide(
    "Projects at Marut Drones",
    "Drone autonomy, computer vision, mapping & GPS-denied navigation",
    "Janadri Yalla Yashwanth",
    "Software Engineer — Drone AI & Automation  |  Technical Interview",
)

plain_slide("Introduction", "About Me", [
    ("Software Engineer at Marut Drones — Drone AI & Automation team", 0),
    ("B.Tech in Computer Science (AI & ML), CGPA 8.6", 0),
    ("Work spans the full drone-software stack:", 0),
    ("Computer vision & ML on drone / aerial imagery", 1),
    ("Autonomous flight, precision landing & obstacle avoidance", 1),
    ("GPS-denied navigation (LiDAR, radar, SLAM) and mapping / photogrammetry", 1),
    ("Production backends that serve the models", 1),
], 2)

PROJECTS = [
    ("UAV Log Analyzer — JIYI & Pixhawk",
     [("Full-stack tool that ingests raw drone logs from two flight controllers, auto-detects the format, decodes, and diagnoses issues", 0),
      ("Dual ecosystem: Pixhawk / ArduPilot (.bin) and JIYI — including XOR-decrypting JIYI's encrypted logs", 1),
      ("From-scratch binary parser (50+ message types); ~30-check engine (power, IMU/vibration, GPS, EKF, compass, RC, failsafes)", 1),
      ("Classifies each finding PASS / WARN / FAIL and infers probable crash cause; generates a PDF report", 1)],
     "Python · FastAPI · React · ReportLab · Docker"),

    ("YOLO Detection Models — Custom Training",
     [("Trained custom YOLO object-detection models on self-labeled datasets for five real-world tasks:", 0),
      ("Water stagnation · broken fence · vehicle detection · PPE-violation detection · fire detection", 1),
      ("Owned the full pipeline: data collection → annotation → augmentation → training → validation & tuning", 1)],
     "YOLO (Ultralytics) · PyTorch · OpenCV"),

    ("Building / Tree / Stockpile Detection — TIFF & Panorama",
     [("Detection pipelines on GeoTIFF orthomosaics and stitched panorama images", 0),
      ("Buildings (LangSAM / GroundingDINO), trees (DeepForest), stockpiles (SAM-2)", 1),
      ("DSM / DTM height filtering to remove false positives; stockpile volume (m3) + material classification", 1),
      ("Output: georeferenced GeoJSON with per-object area, height and volume", 1)],
     "SAM-2 · DeepForest · GroundingDINO · rasterio / GDAL"),

    ("Measurement Tools on GeoTIFF",
     [("Built measurement tools that operate on georeferenced TIFF orthomosaics", 0),
      ("Distance, area and elevation measurement using CRS-aware pixel-to-real-world conversion", 1),
      ("Lets surveyors take accurate real-world measurements directly on the map", 1)],
     "Python · rasterio / GDAL · GeoTIFF"),

    ("Depth Estimation — Stereo Cameras",
     [("Estimates real-world distance to objects from a calibrated stereo camera pair — outputs per-object range (e.g. 2.33 m, 2.17 m)", 0),
      ("Calibrated the cameras with a checkerboard target; stereo rectification → disparity map → metric depth", 1),
      ("Built a custom stereo camera rig running on a Raspberry Pi", 1),
      ("Gives the drone 3D perception / distance sensing without GPS or LiDAR", 1)],
     "OpenCV (stereo vision) · Raspberry Pi · Python"),

    ("Precision Landing on a Moving Platform",
     [("Autonomous landing on an AprilTag target mounted on a moving ground rover", 0),
      ("Raspberry Pi companion: AprilTag (OpenCV) → LANDING_TARGET + velocity setpoints over MAVLink to the flight controller at 20 Hz", 1),
      ("Three-mode state machine: visual control with rover-velocity feed-forward, GPS-NAV, hover fallback; touchdown auto-disarm", 1),
      ("Proven ~5-10 cm accuracy on the stationary tag", 1)],
     "ArduPilot · MAVLink · Raspberry Pi · AprilTag · OpenCV"),

    ("MightySLAM — Research & Hardware Requirements",
     [("Researched MightySLAM for GPS-denied localization and mapping", 0),
      ("Evaluated the approach and its sensor / compute needs for our drone platform", 1),
      ("Produced a hardware-requirements report to guide deployment", 1)],
     "SLAM · Sensor fusion (LiDAR / camera / IMU)"),

    ("Photogrammetry Stitching — WebODM & Pix4D",
     [("Generated orthomosaics and 3D reconstructions from drone image sets", 0),
      ("Used WebODM (open-source) and Pix4D Mapper (commercial) pipelines", 1),
      ("Feeds the stitched maps used by the detection and measurement tools", 1)],
     "WebODM · Pix4D Mapper · Photogrammetry"),

    ("Path Planning & Obstacle Avoidance — Radar",
     [("Pre-mission path planning with radar-based obstacle avoidance, validated on real flight data (226 obstacles mapped in one flight)", 0),
      ("3D-aware avoidance: fly OVER short obstacles, steer around taller ones with a safety margin (2.5 m) and ~3.8 m min clearance", 1),
      ("Compared sensor coverage — a 112° front radar saw 114 of 154 nearby obstacles vs a 360° radar seeing all 154; front + down radar for terrain", 1),
      ("Obstacles and planned routes visualized on satellite imagery", 1)],
     "Radar (360° + front/down) · Path planning · Python · ArduPilot"),

    ("Indoor Navigation without GPS — LiDAR",
     [("Autonomous indoor navigation where GPS is unavailable, using LiDAR-based SLAM", 0),
      ("Built 2D occupancy-grid maps of indoor spaces in real time — walls and free space mapped as the drone moves", 1),
      ("Lets the drone localize and avoid obstacles indoors — key for GPS-denied / EW-resilient operations", 1)],
     "LiDAR · 2D SLAM / occupancy mapping · Python"),

    ("LingBot — 3D Environment Mapping",
     [("Built dense 3D point-cloud reconstructions of the environment — paths, vegetation and structures", 0),
      ("Turns drone / sensor data into a navigable 3D map for perception and planning", 1)],
     "3D reconstruction · Point cloud · SLAM · Python"),
]

# agenda
_agenda = [("%d.  %s" % (i + 1, p[0]), 0) for i, p in enumerate(PROJECTS)]
plain_slide("Agenda", "Projects Overview", _agenda, 3)

# assign numbers + build project slides
slide_no = 4
for i, (t, bullets, tech) in enumerate(PROJECTS, start=1):
    _PROJ_NUM[t] = i
    multi = {
        2: ("2a.jpg", "2b.jpg", "2c.jpg"),
        3: ("3a.jpg", "3b.jpg", "3c.jpg"),
        5: ("5a.jpg", "5b.jpg", "5c.jpg"),
        9: ("9a.png", "9b.png", "9c.png"),
        10: ("10a.jpg", "10b.jpg"),
    }
    if i == 1:
        log_analyzer_slide("Project 1", t, bullets, tech, slide_no)
    elif i in multi:
        imgs = [os.path.join(IMG_DIR, f) for f in multi[i]]
        multi_img_slide("Project %d" % i, t, bullets, tech, slide_no, imgs)
    else:
        project_slide("Project %d" % i, t, bullets, tech, slide_no, img_for(i))
    slide_no += 1

plain_slide("Summary", "Technical Skills", [
    ("Languages: Python, C++, C, Java, SQL", 0),
    ("AI / Computer Vision: PyTorch, YOLO, SAM-2, DeepForest, SegFormer, OpenCV, stereo vision", 0),
    ("Drone / Autonomy: ArduPilot, MAVLink, Raspberry Pi, AprilTag, LiDAR, Radar, SLAM", 0),
    ("Mapping / Geospatial: WebODM, Pix4D, rasterio, GDAL, GeoTIFF, DSM / DTM, GeoJSON", 0),
    ("Backend: FastAPI, Celery, Redis, PostgreSQL, REST APIs", 0),
    ("DevOps: Docker, GitHub Actions, Git, Linux", 0),
], slide_no)
slide_no += 1

section_slide("Thank you  —  Questions?", slide_no)

out = r"c:\Users\yalla\Videos\DISC_14\Marut_Projects_Presentation.pptx"
prs.save(out)
print("SAVED", out, "with", len(prs.slides._sldIdLst), "slides")
